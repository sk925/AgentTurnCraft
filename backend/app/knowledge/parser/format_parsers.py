from __future__ import annotations

from io import BytesIO

from app.knowledge.parser.text_utils import decode_text, split_text_blocks
from app.knowledge.parser.types import DocumentSection


def parse_plain_sections(data: bytes, *, max_block_chars: int) -> list[DocumentSection]:
    text = decode_text(data)
    blocks = split_text_blocks(text, max_chars=max_block_chars)
    return [
        DocumentSection(index=index, content=block, metadata={"block_index": index})
        for index, block in enumerate(blocks)
    ]


def parse_pdf_sections(data: bytes) -> list[DocumentSection]:
    from app.utils.pdf_text import extract_pdf_pages

    sections: list[DocumentSection] = []
    for page in extract_pdf_pages(data):
        sections.append(
            DocumentSection(
                index=len(sections),
                content=page.text,
                metadata={
                    "page_number": page.page_number,
                    "page_count": page.page_count,
                    "text_source": page.source,
                },
            )
        )
    return sections


def parse_docx_sections(data: bytes, *, max_block_chars: int) -> list[DocumentSection]:
    from docx import Document

    doc = Document(BytesIO(data))
    paragraphs = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip()]
    if not paragraphs:
        return []

    text = "\n".join(paragraphs)
    blocks = split_text_blocks(text, max_chars=max_block_chars)
    return [
        DocumentSection(
            index=index,
            content=block,
            metadata={"block_index": index, "paragraph_count": len(paragraphs)},
        )
        for index, block in enumerate(blocks)
    ]


def parse_xlsx_sections(data: bytes, *, rows_per_section: int) -> list[DocumentSection]:
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
    sections: list[DocumentSection] = []
    try:
        for sheet in wb.worksheets:
            batch: list[str] = []
            batch_start_row = 1
            current_row = 0
            for row in sheet.iter_rows(values_only=True):
                current_row += 1
                line = "\t".join("" if cell is None else str(cell) for cell in row).strip()
                if not line:
                    continue
                batch.append(line)
                if len(batch) >= rows_per_section:
                    sections.append(
                        DocumentSection(
                            index=len(sections),
                            content=f"## {sheet.title}\n" + "\n".join(batch),
                            metadata={
                                "sheet_name": sheet.title,
                                "row_start": batch_start_row,
                                "row_end": current_row,
                            },
                        )
                    )
                    batch = []
                    batch_start_row = current_row + 1
            if batch:
                sections.append(
                    DocumentSection(
                        index=len(sections),
                        content=f"## {sheet.title}\n" + "\n".join(batch),
                        metadata={
                            "sheet_name": sheet.title,
                            "row_start": batch_start_row,
                            "row_end": current_row,
                        },
                    )
                )
    finally:
        wb.close()
    return sections


def parse_pptx_sections(data: bytes) -> list[DocumentSection]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(BytesIO(data))
    sections: list[DocumentSection] = []

    def _cell_text(value: str) -> str:
        return " ".join(value.replace("\r", "").split())

    def _collect_shape(shape, lines: list[str], indent: str, *, skip_shape: object | None) -> None:
        if skip_shape is not None and shape is skip_shape:
            return
        shape_type = shape.shape_type
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:  # type: ignore[union-attr]
                _collect_shape(child, lines, indent + "  ", skip_shape=skip_shape)
            return
        if getattr(shape, "has_table", False):
            lines.append(f"{indent}### 表格")
            for row in shape.table.rows:  # type: ignore[union-attr]
                cells = [_cell_text(cell.text) for cell in row.cells]
                lines.append(f"{indent}" + "\t".join(cells))
            lines.append("")
            return
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()  # type: ignore[union-attr]
            if text:
                for block in text.split("\n"):
                    block = block.strip()
                    if block:
                        lines.append(f"{indent}{block}")
            return
        if shape_type in (MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.DIAGRAM, MSO_SHAPE_TYPE.IGX_GRAPHIC):
            labels = {
                MSO_SHAPE_TYPE.CHART: "图表",
                MSO_SHAPE_TYPE.DIAGRAM: "图示",
                MSO_SHAPE_TYPE.IGX_GRAPHIC: "SmartArt",
            }
            lines.append(f"{indent}[{labels.get(shape_type, '嵌入图形')}：无法作为纯文本结构化抽取，已跳过]")

    for slide_number, slide in enumerate(prs.slides, start=1):
        parts: list[str] = [f"## 第 {slide_number} 页"]
        try:
            layout = slide.slide_layout.name.strip()
            if layout:
                parts.append(f"- **版式**: {layout}")
        except (AttributeError, ValueError):
            pass

        title_shape = None
        try:
            title_shape = slide.shapes.title
        except (AttributeError, ValueError):
            pass
        if title_shape is not None and title_shape.text and title_shape.text.strip():
            parts.append(f"- **标题**: {title_shape.text.strip()}")
        parts.append("")

        slide_lines: list[str] = []
        for shape in slide.shapes:
            _collect_shape(shape, slide_lines, "", skip_shape=title_shape)
        chunk = "\n".join(slide_lines).strip()
        if chunk:
            parts.append(chunk)

        content = "\n".join(parts).strip()
        if content:
            sections.append(
                DocumentSection(
                    index=len(sections),
                    content=content,
                    metadata={"slide_number": slide_number, "slide_count": len(prs.slides)},
                )
            )
    return sections
