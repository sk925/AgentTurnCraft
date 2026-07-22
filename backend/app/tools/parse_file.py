from __future__ import annotations

import csv
import json
import re
from html import unescape
from io import BytesIO, StringIO
from pathlib import Path

from bs4 import BeautifulSoup
from langchain.tools import BaseTool, tool
from langchain_core.callbacks import CallbackManagerForToolRun
from minio.error import S3Error
from pydantic import BaseModel, Field, model_validator

from app.chat.base.models.upload_file import UploadFileService
from app.config import settings
from app.enums import FileType
from app.utils.http_fetch import fetch_url_bytes, normalize_http_url
from app.utils.minio_storage import download_bytes

_MAX_CHARS = 100_000

_MIME_TO_TYPE: dict[str, FileType] = {
    "application/pdf": FileType.PDF,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": FileType.DOCX,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": FileType.XLSX,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": FileType.PPTX,
    "text/plain": FileType.TXT,
    "text/csv": FileType.CSV,
    "application/csv": FileType.CSV,
    "application/json": FileType.JSON,
    "application/xml": FileType.XML,
    "text/xml": FileType.XML,
    "text/html": FileType.HTML,
    "text/markdown": FileType.MD,
}


class ParseFileInput(BaseModel):
    """parse_file 工具的入参：file_id 与 url 二选一。"""

    file_id: int | None = Field(
        default=None,
        description="用户上传文件的 ID（对话附件场景）",
    )
    url: str | None = Field(
        default=None,
        description="文件下载链接（http/https），如 PDF、Word、Excel 等直链",
    )

    @model_validator(mode="after")
    def _require_one_source(self) -> ParseFileInput:
        has_id = self.file_id is not None
        has_url = bool(self.url and self.url.strip())
        if has_id == has_url:
            raise ValueError("file_id 与 url 必须且只能提供一个")
        return self


class FileParser(BaseTool):
    name: str = "parse_file"
    description: str = (
        "解析文件内容，支持 txt/md/csv/json/html/xml/docx/xlsx/pptx/pdf；"
        "可通过 file_id 读取用户上传的文件，或通过 url 下载并解析远程文件（含 PDF 公告等）；"
        "扫描版 PDF 会在文字层为空时自动 OCR 识别。"
        "查看普通网页正文请优先使用 fetch_webpage。"
    )
    args_schema: type[BaseModel] = ParseFileInput

    def _run(
        self,
        file_id: int | None = None,
        url: str | None = None,
        run_manager: CallbackManagerForToolRun | None = None,
    ) -> str:
        return parse_file_content(file_id=file_id, url=url)


def _resolve_file_type(mime_or_ext: str, file_name: str) -> FileType | None:
    """结合 MIME（或扩展名）与文件名推断 FileType。"""
    raw = (mime_or_ext or "").strip().lower()
    if raw in {e.value for e in FileType}:
        return FileType(raw)
    base = raw.split(";", 1)[0].strip()
    if base in _MIME_TO_TYPE:
        return _MIME_TO_TYPE[base]
    if raw in _MIME_TO_TYPE:
        return _MIME_TO_TYPE[raw]
    if base == "application/x-pdf":
        return FileType.PDF

    ext = Path(file_name).suffix.lower().lstrip(".")
    if not ext:
        return None
    try:
        return FileType(ext)
    except ValueError:
        return None


def _decode_text(data: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "gbk", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _truncate(text: str) -> str:
    if len(text) <= _MAX_CHARS:
        return text
    return (
        text[:_MAX_CHARS]
        + f"\n\n[内容过长已截断，共 {len(text)} 字符，仅保留前 {_MAX_CHARS} 字符]"
    )


def _parse_plain(data: bytes) -> str:
    return _decode_text(data)


def _parse_html(data: bytes) -> str:
    html = _decode_text(data)
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript", "svg", "iframe"]):
        tag.decompose()
    text = soup.get_text("\n", strip=True)
    return unescape(re.sub(r"\n{3,}", "\n\n", text))


def _parse_json_text(data: bytes) -> str:
    text = _decode_text(data)
    try:
        obj = json.loads(text)
        return json.dumps(obj, ensure_ascii=False, indent=2)
    except json.JSONDecodeError:
        return text


def _parse_csv(data: bytes) -> str:
    text = _decode_text(data)
    buf = StringIO(text)
    rows = list(csv.reader(buf))
    return "\n".join("\t".join(row) for row in rows)


def _parse_pdf(data: bytes) -> str:
    from app.utils.pdf_text import extract_pdf_full_text

    return extract_pdf_full_text(data)


def _parse_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _parse_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
    try:
        blocks: list[str] = []
        for sheet in wb.worksheets:
            blocks.append(f"## {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                line = "\t".join("" if c is None else str(c) for c in row)
                blocks.append(line)
        return "\n".join(blocks)
    finally:
        wb.close()


def _parse_pptx(data: bytes) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(BytesIO(data))
    parts: list[str] = []

    def _cell_text(s: str) -> str:
        return " ".join(s.replace("\r", "").split())

    def _collect_shape(shape, lines: list[str], indent: str, *, skip_shape: object | None) -> None:
        if skip_shape is not None and shape is skip_shape:
            return
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:  # type: ignore[union-attr]
                _collect_shape(child, lines, indent + "  ", skip_shape=skip_shape)
            return
        if getattr(shape, "has_table", False):
            lines.append(f"{indent}### 表格")
            for row in shape.table.rows:  # type: ignore[union-attr]
                cells = [_cell_text(c.text) for c in row.cells]
                lines.append(f"{indent}" + "\t".join(cells))
            lines.append("")
            return
        if getattr(shape, "has_text_frame", False):
            t = shape.text_frame.text.strip()  # type: ignore[union-attr]
            if t:
                for block in t.split("\n"):
                    block = block.strip()
                    if block:
                        lines.append(f"{indent}{block}")
            return
        if st in (MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.DIAGRAM, MSO_SHAPE_TYPE.IGX_GRAPHIC):
            kind_labels = {
                MSO_SHAPE_TYPE.CHART: "图表",
                MSO_SHAPE_TYPE.DIAGRAM: "图示",
                MSO_SHAPE_TYPE.IGX_GRAPHIC: "SmartArt",
            }
            label = kind_labels.get(st, "嵌入图形")
            lines.append(f"{indent}[{label}：无法作为纯文本结构化抽取，已跳过]")

    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"## 第 {idx} 页")
        try:
            layout = slide.slide_layout.name.strip()
            if layout:
                parts.append(f"- **版式**: {layout}")
        except (AttributeError, ValueError):
            pass

        title_ph = None
        try:
            title_ph = slide.shapes.title
        except (AttributeError, ValueError):
            pass
        if title_ph is not None and title_ph.text and title_ph.text.strip():
            parts.append(f"- **标题**: {title_ph.text.strip()}")
        parts.append("")

        slide_lines: list[str] = []
        for shape in slide.shapes:
            _collect_shape(shape, slide_lines, "", skip_shape=title_ph)

        chunk = "\n".join(slide_lines).strip()
        if chunk:
            parts.append(chunk)
        parts.append("")

    return "\n".join(parts).strip()


def _parse_by_type(kind: FileType, data: bytes) -> str:
    if kind == FileType.HTML:
        return _parse_html(data)
    if kind in (FileType.TXT, FileType.MD, FileType.XML):
        return _parse_plain(data)
    if kind == FileType.JSON:
        return _parse_json_text(data)
    if kind == FileType.CSV:
        return _parse_csv(data)
    if kind == FileType.PDF:
        return _parse_pdf(data)
    if kind == FileType.DOCX:
        return _parse_docx(data)
    if kind == FileType.XLSX:
        return _parse_xlsx(data)
    if kind == FileType.PPTX:
        return _parse_pptx(data)
    raise ValueError(f"未实现的类型: {kind}")


def _parse_bytes(data: bytes, *, mime_or_ext: str, file_name: str, source: str) -> str:
    if not data:
        return "文件内容为空"

    kind = _resolve_file_type(mime_or_ext, file_name)
    if kind is None:
        allowed = ", ".join(e.value for e in FileType)
        return (
            f"无法识别文件类型（type={mime_or_ext!r}, 文件名={file_name!r}）。"
            f"支持的类型：{allowed}\n来源: {source}"
        )

    try:
        text = _parse_by_type(kind, data).strip()
    except Exception as e:
        return f"解析失败: {e}\n来源: {source}"

    if not text:
        return "解析结果为空（可能为加密文档、纯图片文件，或 OCR 未能识别到文字）"

    header = f"来源: {source}\n类型: {kind.value}\n\n"
    return _truncate(header + text)


def parse_file_by_id(file_id: int) -> str:
    """根据文件 ID 解析用户上传的文件内容。"""
    upload_file = UploadFileService.get_upload_file_by_id(file_id)
    if not upload_file:
        return "文件不存在"

    try:
        data = download_bytes(settings.minio_bucket, upload_file.file_path)
    except S3Error as e:
        return f"从对象存储读取文件失败: {e.message}"

    return _parse_bytes(
        data,
        mime_or_ext=upload_file.file_type,
        file_name=upload_file.file_name,
        source=f"file_id={file_id} ({upload_file.file_name})",
    )


def parse_file_by_url(url: str) -> str:
    """根据 URL 下载并解析远程文件内容。"""
    try:
        normalized = normalize_http_url(url)
        fetched = fetch_url_bytes(normalized)
    except ValueError as exc:
        return f"无效请求: {exc}"
    except Exception as exc:
        return f"下载失败: {exc}"

    return _parse_bytes(
        fetched.data,
        mime_or_ext=fetched.content_type,
        file_name=fetched.filename_hint,
        source=fetched.final_url,
    )


def parse_file_content(*, file_id: int | None = None, url: str | None = None) -> str:
    """统一入口：按 file_id 或 url 解析文件。"""
    if file_id is not None:
        return parse_file_by_id(file_id)
    if url:
        return parse_file_by_url(url)
    return "请提供 file_id 或 url 之一"


@tool(
    "parse_file_by_id",
    description="根据用户上传文件的 ID 解析文件内容（仅 file_id，不含 URL）。",
)
def parse_file_by_id_tool(file_id: int) -> str:
    """供 agent_selector 等场景使用的 file_id 专用工具。"""
    return parse_file_by_id(file_id)
